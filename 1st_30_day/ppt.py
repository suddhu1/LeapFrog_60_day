# Create a PowerPoint presentation with Object Design Process slides and a simple UML diagram

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize presentation
prs = Presentation()

# Helper function to add a slide with title and bullet points
def add_bullet_slide(title, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, point in enumerate(bullet_points):
        if i == 0:
            body.text = point
        else:
            body.add_paragraph().text = point

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Object Design Process"
slide.placeholders[1].text = "An Overview\n\nPresented by: Your Name"

# Slide 2: Introduction
add_bullet_slide(
    "Introduction",
    [
        "Object Design is a phase of Object-Oriented Software Engineering",
        "It follows Object-Oriented Analysis (OOA)",
        "Focuses on how the system will be built",
        "Transforms analysis into design models"
    ]
)

# Slide 3: What is Object Design?
add_bullet_slide(
    "What is Object Design?",
    [
        "Defines classes, attributes, and methods",
        "Establishes relationships between objects",
        "Acts as a bridge between analysis and coding",
        "Produces detailed design specifications"
    ]
)

# Slide 4: Objectives of Object Design
add_bullet_slide(
    "Objectives of Object Design",
    [
        "Create a clear system structure",
        "Assign responsibilities to objects",
        "Improve reusability and maintainability",
        "Reduce system complexity"
    ]
)

# Slide 5: Key Activities
add_bullet_slide(
    "Key Activities in Object Design",
    [
        "Identifying classes",
        "Defining attributes",
        "Defining methods",
        "Establishing relationships",
        "Applying design principles"
    ]
)

# Slide 6: Identifying Classes
add_bullet_slide(
    "Identifying Classes",
    [
        "Derived from real-world entities",
        "Extracted from use cases",
        "Each class represents an object",
        "Example: Student, Course, Teacher"
    ]
)

# Slide 7: Defining Attributes
add_bullet_slide(
    "Defining Attributes",
    [
        "Attributes define the state of an object",
        "Should be relevant and necessary",
        "Avoid redundancy",
        "Example: Student → name, rollNo, age"
    ]
)

# Slide 8: Defining Methods
add_bullet_slide(
    "Defining Methods",
    [
        "Methods define object behavior",
        "Perform operations on attributes",
        "Support interaction between objects",
        "Example: enrollCourse(), displayInfo()"
    ]
)

# Slide 9: Object Relationships
add_bullet_slide(
    "Object Relationships",
    [
        "Objects communicate using relationships",
        "Association, Aggregation, Composition",
        "Inheritance for reusability",
        "Example: Student enrolls in Course"
    ]
)

# Slide 10: Design Principles
add_bullet_slide(
    "Design Principles",
    [
        "Encapsulation",
        "Abstraction",
        "Low coupling and high cohesion",
        "Modularity and flexibility"
    ]
)

# Slide 11: UML Diagrams
add_bullet_slide(
    "UML Diagrams in Object Design",
    [
        "Class Diagram",
        "Sequence Diagram",
        "Object Diagram",
        "Used to visualize system structure"
    ]
)

# Slide 12: UML Class Diagram (Visual)
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Sample UML Class Diagram"

# Student class box
student_box = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), Inches(3), Inches(2)
)
student_box.text_frame.text = "Student\n----------------\nname\nrollNo\nage\n----------------\nenrollCourse()\ndisplayInfo()"

# Course class box
course_box = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(5), Inches(1.5), Inches(3), Inches(2)
)
course_box.text_frame.text = "Course\n----------------\ncourseName\ncourseId\n----------------\naddStudent()"

# Slide 13: Benefits
add_bullet_slide(
    "Benefits of Object Design",
    [
        "Clear architecture",
        "Easy implementation",
        "Scalable system",
        "Better code quality"
    ]
)

# Slide 14: Conclusion
add_bullet_slide(
    "Conclusion",
    [
        "Object Design converts requirements into structure",
        "Essential for robust object-oriented systems",
        "Good design leads to quality software"
    ]
)

# Slide 15: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Thank You"
slide.placeholders[1].text = "Questions?"

# Save presentation
file_path = "/home/sudarshan/Desktop/LeapFrog_60_day/Object_Design_Process_Presentation.pptx"
prs.save(file_path)

file_path
